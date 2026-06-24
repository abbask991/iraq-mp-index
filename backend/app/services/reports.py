"""Server-side PDF reports via Playwright (chromium).

For serious, government-ready documents — executive reports, campaign dossiers,
political profiles. html2pdf.js stays for quick client-side exports; this path
produces consistent, print-quality A4 PDFs rendered headless.

Playwright is lazy-imported so the web process doesn't need chromium — only the
worker (which runs `generate_report`) installs it. Without it, build() returns a
clear error instead of crashing.
"""
import base64
import html as _html


def _esc(s):
    return _html.escape(str(s if s is not None else ""))


def _shell(title, subtitle, body):
    """Wrap report body in a branded, RTL, print-optimized A4 document."""
    return f"""<!doctype html><html dir="rtl" lang="ar"><head><meta charset="utf-8">
<style>
  @page {{ size: A4; margin: 18mm 14mm; }}
  * {{ font-family: 'Segoe UI', Tahoma, Arial, sans-serif; box-sizing: border-box; }}
  body {{ color: #1a2230; margin: 0; }}
  .head {{ border-bottom: 3px solid #0b5; padding-bottom: 12px; margin-bottom: 18px; }}
  .brand {{ color: #0b5; font-weight: 800; font-size: 13px; letter-spacing: .5px; }}
  h1 {{ font-size: 22px; margin: 6px 0 2px; }}
  .sub {{ color: #5a6677; font-size: 12px; }}
  h2 {{ font-size: 15px; margin: 20px 0 8px; color: #0a3; border-right: 4px solid #0b5; padding-right: 8px; }}
  .kpis {{ display: flex; gap: 10px; flex-wrap: wrap; margin: 10px 0; }}
  .kpi {{ flex: 1; min-width: 110px; background: #f4f8f5; border: 1px solid #e2ece6; border-radius: 8px; padding: 10px; text-align: center; }}
  .kpi b {{ display: block; font-size: 20px; color: #084; }}
  .kpi span {{ font-size: 11px; color: #5a6677; }}
  table {{ width: 100%; border-collapse: collapse; font-size: 12px; margin: 6px 0; }}
  th, td {{ border: 1px solid #e2ece6; padding: 6px 8px; text-align: right; }}
  th {{ background: #f0f6f2; }}
  p {{ font-size: 13px; line-height: 1.7; }}
  .foot {{ margin-top: 26px; border-top: 1px solid #e2ece6; padding-top: 8px; color: #8a93a0; font-size: 10px; }}
  .tag {{ display: inline-block; background: #eef4f0; border-radius: 6px; padding: 2px 8px; font-size: 11px; margin: 2px; }}
</style></head><body>
  <div class="head"><div class="brand">مركز الرصد · RASD INTELLIGENCE</div>
    <h1>{_esc(title)}</h1><div class="sub">{_esc(subtitle)}</div></div>
  {body}
  <div class="foot">تقرير آلي صادر عن منصة مركز الرصد — المؤشرات احتمالية وتتطلّب مراجعة بشرية قبل أي قرار.
  وثيقة سرّية لأغراض التحليل.</div>
</body></html>"""


def _kpis(pairs):
    cells = "".join(f'<div class="kpi"><b>{_esc(v)}</b><span>{_esc(k)}</span></div>' for k, v in pairs)
    return f'<div class="kpis">{cells}</div>'


def _table(headers, rows):
    head = "".join(f"<th>{_esc(h)}</th>" for h in headers)
    body = "".join("<tr>" + "".join(f"<td>{_esc(c)}</td>" for c in r) + "</tr>" for r in rows)
    return f"<table><tr>{head}</tr>{body}</table>"


# ---- data → HTML per report kind ----
def _profile_html(d):
    s = d.get("sentiment", {})
    body = _kpis([("إجمالي المنشورات", d.get("total", 0)),
                  ("المؤشر الإعلامي", f"{d.get('media_index', '-')}/100"),
                  ("إيجابي", s.get("pos", 0)), ("سلبي", s.get("neg", 0))])
    if d.get("executive"):
        body += f"<h2>التقييم التنفيذي</h2><p>{_esc(d['executive'])}</p>"
    if d.get("content", {}).get("narratives"):
        body += "<h2>السرديات المهيمنة</h2>" + _table(
            ["السردية", "الوصف", "الحصة %"],
            [[n.get("label"), n.get("description"), n.get("share")] for n in d["content"]["narratives"]])
    if d.get("campaign", {}).get("score") is not None:
        c = d["campaign"]
        body += f"<h2>إشارات التنسيق</h2><p>درجة التنسيق: <b>{_esc(c.get('score'))}</b> — {_esc(c.get('level'))}</p>"
    if d.get("top_items"):
        body += "<h2>أبرز المنشورات</h2>" + _table(
            ["المصدر", "النبرة", "العنوان"],
            [[i.get("source"), i.get("sentiment"), (i.get("title") or "")[:90]] for i in d["top_items"][:8]])
    return body


def _campaign_html(d):
    body = _kpis([("درجة التنسيق", d.get("coordination_score", 0)),
                  ("المنشورات", d.get("total_posts", 0)),
                  ("الحسابات", d.get("unique_accounts", 0)),
                  ("التصنيف", (d.get("alert_level") or {}).get("label", "-"))])
    body += f"<h2>التفسير</h2><p>{_esc(d.get('explanation', ''))}</p>"
    if d.get("sub_scores"):
        body += "<h2>الإشارات التسعة</h2>" + _table(
            ["الإشارة", "الدرجة"], [[k, v] for k, v in d["sub_scores"].items()])
    if d.get("top_hashtags"):
        body += "<h2>الهاشتاغات</h2>" + "".join(
            f'<span class="tag">{_esc(h["hashtag"])} ({h["count"]})</span>' for h in d["top_hashtags"][:12])
    return body


def _executive_html(d):
    s = d.get("sentiment", {})
    body = _kpis([("المسوح", d.get("scanned", d.get("total", 0))),
                  ("الحسابات", d.get("accounts", 0)),
                  ("المؤشر", f"{d.get('media_index', '-')}/100"),
                  ("سلبي", s.get("neg", 0))])
    if d.get("trending"):
        body += "<h2>الترندات</h2>" + "".join(
            f'<span class="tag">{_esc(h.get("hashtag", h))}</span>' for h in d["trending"][:12])
    if d.get("campaigns"):
        body += "<h2>حملات مشتبهة</h2>" + _table(
            ["الهاشتاغ", "درجة التنسيق"], [[c.get("hashtag"), c.get("coordination_score")] for c in d["campaigns"][:8]])
    return body


async def _gather(kind, target, rng):
    """Pull the data payload for a report kind from the existing endpoints."""
    from app.routers import monitor as m
    if kind == "campaign":
        return await m.monitor_campaign(m.KeywordReq(keywords=[target], range=rng))
    if kind == "executive":
        return await m.monitor_overview(m.KeywordReq(range=rng))
    # profile / government default → full dossier
    return await m.monitor_dossier(m.KeywordReq(keywords=[target], range=rng))


def _render_html(kind, target, data):
    title = {"campaign": "تقرير حملة منظّمة", "executive": "التقرير التنفيذي للرصد",
             "government": f"ملف استخباراتي — {target}"}.get(kind, f"الملف الشامل — {target}")
    sub = f"النطاق الزمني: {data.get('period', data.get('window', '—'))} · تاريخ الإصدار آلي"
    builder = {"campaign": _campaign_html, "executive": _executive_html}.get(kind, _profile_html)
    return _shell(title, sub, builder(data))


# ---- neutral document model (shared by docx / pptx) ----
def _doc_model(kind, target, data):
    title = {"campaign": "تقرير حملة منظّمة", "executive": "التقرير التنفيذي للرصد",
             "government": f"ملف استخباراتي — {target}"}.get(kind, f"الملف الشامل — {target}")
    sub = f"النطاق الزمني: {data.get('period', data.get('window', '—'))} · إصدار آلي"
    s = data.get("sentiment", {})
    kpis = [("إجمالي", data.get("total", data.get("scanned", data.get("total_posts", 0)))),
            ("المؤشر الإعلامي", f"{data.get('media_index', '-')}/100"),
            ("إيجابي", s.get("pos", 0)), ("سلبي", s.get("neg", 0))]
    sections = []
    if data.get("executive"):
        sections.append(("التقييم التنفيذي", [data["executive"]]))
    if data.get("explanation"):
        sections.append(("التفسير", [data["explanation"]]))
    narrs = data.get("content", {}).get("narratives") or data.get("narratives") or []
    if narrs:
        sections.append(("السرديات المهيمنة",
                         [f"• {n.get('label')}: {n.get('description','')} ({n.get('share','')}%)" for n in narrs[:5]]))
    if data.get("sub_scores"):
        sections.append(("إشارات التنسيق",
                         [f"• {k}: {v}" for k, v in data["sub_scores"].items()]))
    return {"title": title, "subtitle": sub, "kpis": kpis, "sections": sections}


def _build_docx(kind, target, data):
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except Exception:
        return {"error": "docx_unavailable", "message": "python-docx runs on the worker."}
    m = _doc_model(kind, target, data)
    doc = Document()
    h = doc.add_heading(m["title"], level=0)
    doc.add_paragraph(m["subtitle"])
    t = doc.add_table(rows=1, cols=len(m["kpis"]))
    for i, (k, v) in enumerate(m["kpis"]):
        c = t.rows[0].cells[i]
        c.text = f"{v}\n{k}"
    for heading, paras in m["sections"]:
        doc.add_heading(heading, level=1)
        for p in paras:
            doc.add_paragraph(str(p))
    doc.add_paragraph("تقرير آلي — المؤشرات احتمالية وتتطلّب مراجعة بشرية. وثيقة سرّية.")
    import io
    buf = io.BytesIO(); doc.save(buf)
    return {"kind": kind, "target": target, "format": "docx", "bytes": buf.tell(),
            "file_base64": base64.b64encode(buf.getvalue()).decode("ascii")}


def _build_pptx(kind, target, data):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return {"error": "pptx_unavailable", "message": "python-pptx runs on the worker."}
    m = _doc_model(kind, target, data)
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = m["title"]
    title_slide.placeholders[1].text = m["subtitle"]
    kpi_slide = prs.slides.add_slide(prs.slide_layouts[1])
    kpi_slide.shapes.title.text = "المؤشرات الرئيسية"
    body = kpi_slide.placeholders[1].text_frame
    body.text = " | ".join(f"{k}: {v}" for k, v in m["kpis"])
    for heading, paras in m["sections"]:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = heading
        tf = sl.placeholders[1].text_frame
        tf.text = str(paras[0]) if paras else ""
        for p in paras[1:]:
            tf.add_paragraph().text = str(p)
    import io
    buf = io.BytesIO(); prs.save(buf)
    return {"kind": kind, "target": target, "format": "pptx", "bytes": buf.tell(),
            "file_base64": base64.b64encode(buf.getvalue()).decode("ascii")}


async def build(kind: str, target: str, rng: str = "week", fmt: str = "pdf") -> dict:
    """Gather data → render report in `fmt` (pdf | docx | pptx). Returns base64."""
    data = await _gather(kind, target, rng)
    if fmt == "docx":
        return _build_docx(kind, target, data)
    if fmt == "pptx":
        return _build_pptx(kind, target, data)
    doc = _render_html(kind, target, data)
    try:
        from playwright.async_api import async_playwright
    except Exception:
        return {"error": "playwright_unavailable",
                "message": "PDF rendering runs on the worker (playwright + chromium).",
                "html": doc}
    async with async_playwright() as p:
        browser = await p.chromium.launch(args=["--no-sandbox"])
        page = await browser.new_page()
        await page.set_content(doc, wait_until="networkidle")
        pdf = await page.pdf(format="A4", print_background=True)
        await browser.close()
    return {"kind": kind, "target": target, "format": "pdf", "bytes": len(pdf),
            "pdf_base64": base64.b64encode(pdf).decode("ascii")}
